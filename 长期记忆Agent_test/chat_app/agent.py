import hashlib
import json
import logging
import os
import sqlite3
from typing import Optional, List, Dict

import chromadb
import numpy as np
from chromadb import Settings
from chromadb.utils import embedding_functions
from openai import OpenAI
from docx import Document
from pptx import Presentation

class LongTermMemoryAgent:
    def __init__(self, db_path="Red_memory3.db",chroma_path = "./chroma.db"):
        self.client = OpenAI(api_key="sk-X9xETgX9p8lZboXMvWQTy7mHKcpmBO1vTJ4sNzsNIB2YohyG",
                             base_url="https://api.chatanywhere.tech/v1")
        self.db_path = db_path
        # 配置参数
        # 向量数据库配置
        # 改用PersistentClient初始化（适配0.4.0+版本）
        self.chroma_client = chromadb.PersistentClient(
            path=chroma_path,  # 替代persist_directory
            settings=Settings(anonymized_telemetry=False)
        )
        self.embedding_model = "text-embedding-3-small"
        self.chat_model = "gpt-3.5-turbo-1106"
        self.memory_decay = 0.98  # 每日记忆衰减率
        self.top_k_memories = 5  # 每次检索记忆数量
        self.ef = embedding_functions.OpenAIEmbeddingFunction(
            api_key="sk-X9xETgX9p8lZboXMvWQTy7mHKcpmBO1vTJ4sNzsNIB2YohyG",
            model_name=self.embedding_model,
            api_base="https://api.chatanywhere.tech/v1"
        )
        # 初始化集合
        self.collection = self.chroma_client.get_or_create_collection(
            name="memories",
            embedding_function=self.ef
        )


    def is_username_exist(self, username) -> bool:
        with sqlite3.connect(self.db_path) as conn:
            conn.row_factory = sqlite3.Row
            cursor = conn.cursor()
            if username:
                cursor.execute("SELECT * FROM chat_app_user WHERE username=?", (username,))
            result = cursor.fetchone()
            return True if result else False

    # ✅新写的
    def get_user_password(self, username):
        with sqlite3.connect(self.db_path) as conn:
            cursor=conn.cursor()
            cursor.execute("SELECT password FROM chat_app_user WHERE username=?",(username,))
            result = cursor.fetchone()
            return result[0] if result else None

    def add_user(self, username: str) -> int:
        """添加用户，返回的是用户的id"""
        with sqlite3.connect(self.db_path) as conn:
            cursor = conn.cursor()
            cursor.execute("""INSERT OR REPLACE INTO users(username)VALUES(?)""", (username,))
            conn.commit()
            return cursor.lastrowid

    def search_user(self, username: str) -> Optional[int]:
        """根据用户名字查找用户表的user_id，再通过查找user_id去session表里的会话"""
        with sqlite3.connect(self.db_path) as conn:
            cursor = conn.cursor()
            cursor.execute("SELECT id FROM chat_app_user WHERE username=?", (username,))
            result = cursor.fetchone()
            return result[0] if result else None

    def creat_session(self, user_id: int, user_input: str):
        """根据提示的summary还有些别的来创建会话表"""
        # 先是调用chat函数，再是将提问总结成summary，metadata，status，保存到session
        summary_response = self.client.chat.completions.create(
            model=self.chat_model,
            messages=[{
                "role": "system", "content": "你是一个信息总结处理的LLM助手。\n"
                                             "需要你把用户提问的信息总结成一句不超过10个字的主题\n"
            }, {"role": "user", "content": f"用户说{user_input}"}
            ]
        )
        summary = summary_response.choices[0].message.content
        metadata_response = self.client.chat.completions.create(
            model=self.chat_model,
            messages=[{
                "role": "system", "content": "你是一个信息提取处理的LLM助手。\n"
                                             "需要你对用户的信息提取出一个metadata，即会话的主题以及用户的提示词。"
            }, {"role": "user", "content": f"用户说{user_input}"}
            ]
        )
        metadata = metadata_response.choices[0].message.content
        status = "active"
        with sqlite3.connect(self.db_path) as conn:
            cursor = conn.cursor()
            cursor.execute("""
                INSERT  INTO chat_app_chatsession
                (user_id,last_accessed,summary,status,metadata)
                VALUES(?,CURRENT_TIMESTAMP,?,?,?)""", (user_id, summary, status, str(metadata)))
            session_id = cursor.lastrowid  # 获取新生成的 session_id
            conn.commit()
        return session_id

    def export_sessions(self, user_id: int) :
        with sqlite3.connect(self.db_path) as conn:
            conn.row_factory = sqlite3.Row
            cursor = conn.cursor()
            cursor.execute("SELECT summary FROM chat_app_chatsession WHERE user_id = ?", (user_id,))
            sessions = cursor.fetchall()
            if not sessions:
                return None
            return sessions
            # for row in sessions:
            #     d = dict(row)
            #     print(d.values(), type(str))

    def export_session_messages(self, user_input_summary: str) -> Optional[List[Dict]]:
        """感觉还是需要先后存user_input,再是ai_response，否则无法按照时间来陈列出来"""
        with sqlite3.connect(self.db_path) as conn:
            conn.row_factory = sqlite3.Row
            cursor = conn.cursor()
            cursor.execute("""
                SELECT role,message,last_accessed
                FROM chat_app_chatmessage
                WHERE session_id=(SELECT id FROM chat_app_chatsession WHERE summary=?)""", (user_input_summary,))
            messages = cursor.fetchall()
            return messages
            # for message in messages:
            #     mes = dict(message)
            #     print(mes.values(), type(str))
            # return [dict(row)for row in message] if message else None

    @staticmethod
    def _cosine_similarity(embedding1: bytes, embedding2: bytes) -> float:
        """计算余弦相似度（sqlite自定义函数）"""
        vec1 = np.frombuffer(embedding1, dtype=np.float32)
        vec2 = np.frombuffer(embedding2, dtype=np.float32)
        return float(np.dot(vec1, vec2) / (np.linalg.norm(vec1) * np.linalg.norm(vec2)))

    def _get_memory(self, text: str) -> bytes:
        """获取文本记忆，并将数据序列化"""
        response = self.client.embeddings.create(
            input=text,
            model=self.embedding_model
        )
        vector = response.data[0].embedding  # 化为1536维度向量
        return np.array(vector).astype(np.float32).tobytes()  # 精度压缩减少内存并转化成字节数据输出

    def add_memory(self, text: str, role: str, user_id: int, metadata: dict = None,file_type: str = None):
        # message_id = hashlib.md5(text.encode()).hexdigest()
        message_id = hashlib.md5(text.encode()).hexdigest()
        embedding = self._get_memory(text)
        try:
            with sqlite3.connect(self.db_path) as conn:
                session_id = conn.execute('''
                    SELECT id 
                    FROM chat_app_chatsession
                    WHERE user_id=?
                    ORDER BY last_accessed DESC 
                    LIMIT 1
                    ''', (user_id,)).fetchone()
                if not session_id:
                    self.creat_session(user_id, text)
                    session_id = conn.execute('SELECT last_insert_rowid()').fetchone()[0]
                else:
                    session_id = session_id[0]
                # 准备元数据
                entry_metadata = {
                    "role": role,
                    "user_id": user_id,
                    "session_id": session_id,
                    "timestamp": json.dumps(sqlite3.datetime.datetime.now().isoformat()),
                    "weight": 1.0,
                    "file_type": file_type or "text",
                    "is_office_file": file_type in ["docx", "doc", "pptx","ppt"]  # 新增：标记是否为Office文件
                }
                if metadata:
                    entry_metadata.update(metadata)
                emb_resp = self.client.embeddings.create(input=text, model="text-embedding-3-small")
                emb_vec = np.array(emb_resp.data[0].embedding, dtype=np.float32)
                emb_bytes = emb_vec.tobytes()

                # 添加到向量数据库
                self.collection.upsert(
                    ids=[message_id],
                    documents=[text],
                    metadatas=[entry_metadata],
                    embeddings=[emb_vec.tolist()]
                )
                cursor = conn.cursor()

                cursor.execute('''
                    INSERT OR REPLACE INTO chat_app_chatmessage
                    (message_id,session_id,user_id,role,message,last_accessed,embedding,weight,metadata)
                    VALUES (?,?,?,?,?, CURRENT_TIMESTAMP,?,
                    COALESCE((SELECT weight FROM chat_app_chatmessage WHERE message_id =?)*1.1,1.0),?)
                    ''', (message_id, session_id, user_id, role, text, embedding, message_id, str(metadata)))
                conn.commit()
        except Exception as e:
            logging.warning(e)


    def apply_memory_decay(self):
        with sqlite3.connect(self.db_path) as conn:
            conn.execute('''
                UPDATE chat_app_chatmessage
                SET weight=weight*?
                WHERE DATE (last_accessed) <DATE('now','-1 day')
                ''', (self.memory_decay,))
            conn.execute('''
                DELETE FROM chat_app_chatmessage
                WHERE weight<0.1
                ''')

    # def retrieve_messages(self, query: str, session_id: int = None, role: str = None) -> list:
    #     """
    # 根据query语义和session_id检索指定会话的历史消息
    # :param query: 搜索文本
    # :param session_id: 会话ID（需明确传入）
    # :param user_id: 用户ID（用于校验权限）
    # :param role_filter: 角色过滤（'user','assistant','system'）
    # :return: 匹配消息列表，格式：[{"role": "user", "message": "..."}, ...]
    # """
    #     # validate_session_ownership(session_id,)
    #     self.apply_memory_decay()
    #     query_embedding = self._get_memory(query)
    #     with sqlite3.connect(self.db_path) as conn:
    #         conn.create_function("cosine_similarity", 2, self._cosine_similarity)
    #         conditions = ["session_id=?"]
    #         params = [session_id]
    #         if role:
    #             conditions.append("role=?")
    #             params.append(role)
    #         cursor = conn.execute(f'''
    #             SELECT role,message,metadata
    #             FROM chat_app_chatmessage
    #             WHERE {' AND '.join(conditions)}
    #             ORDER BY cosine_similarity(embedding,?)*weight DESC
    #             LIMIT ?
    #                 ''', params + [query_embedding, self.top_k_memories])
    #         messages = cursor.fetchall()
    #         return [{
    #             "role": m[0],
    #             "message": m[1],
    #             "metadata": m[2] if m[2] else None
    #         } for m in messages]

    def retrieve_messages(self, query: str, session_id: int = None, role: str = None) -> list:
        """从向量数据库检索相关消息（兼容Chroma旧版本的单条件查询）"""
        self.apply_memory_decay()

        # 构建过滤条件：只保留一个核心条件（优先session_id，因为会话ID过滤更精准）
        where_clause = None
        if session_id:
            where_clause = {"session_id": session_id}  # 仅保留会话ID作为过滤条件
        elif role:
            where_clause = {"role": role}  # 无会话ID时，再用角色过滤

        # 检索相似文档（最多返回top_k_memories*2条，留足筛选空间）
        results = self.collection.query(
            query_texts=[query],
            n_results=self.top_k_memories * 2,  # 扩大检索范围
            where=where_clause  # 仅传入单个条件（兼容旧版本）
        )

        # 手动过滤第二个条件（role或session_id）
        filtered_docs = []
        filtered_metas = []
        for doc, meta in zip(results['documents'][0], results['metadatas'][0]):
            # 检查是否满足所有条件
            match = True
            if session_id and str(meta.get("session_id")) != str(session_id):  # 注意类型一致（可能是字符串）
                match = False
            if role and meta.get("role") != role:
                match = False
            if match:
                filtered_docs.append(doc)
                filtered_metas.append(meta)

        # 格式化结果（只保留前top_k_memories条）
        messages = []
        for doc, meta in zip(filtered_docs[:self.top_k_memories], filtered_metas[:self.top_k_memories]):
            messages.append({
                "role": meta["role"],
                "message": doc,
                "metadata": meta
            })
        return messages

    def process_file(self, file_path: str, file_type: str, user_id: int) -> str:
        """处理上传的文件并生成摘要（新增Word/PPT解析）"""
        content = ""
        file_ext = os.path.splitext(file_path)[1].lower()  # 获取文件后缀（.docx/.pptx等）

        # 1. 文本文件
        if file_type.startswith('text') or file_ext in ['.txt', '.csv']:
            with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                content = f.read()

        # 2. 图片文件
        elif file_type.startswith('image') or file_ext in ['.jpg', '.jpeg', '.png', '.gif']:
            import base64
            with open(file_path, "rb") as f:
                base64_image = base64.b64encode(f.read()).decode('utf-8')

            response = self.client.chat.completions.create(
                model="gpt-4-vision-preview",
                messages=[
                    {
                        "role": "user",
                        "content": [
                            {"type": "text", "text": "请描述这张图片的内容并提取关键信息"},
                            {"type": "image_url", "image_url": {"url": f"data:image/jpeg;base64,{base64_image}"}}
                        ]
                    }
                ],
                max_tokens=300
            )
            content = response.choices[0].message.content

        # 3. Word文件（.docx/.doc）
        elif file_ext in ['.docx', '.doc']:
            try:
                if file_ext == '.docx':
                    # 解析.docx（python-docx）

                    doc = Document(file_path)
                    content = "\n".join([para.text for para in doc.paragraphs if para.text.strip()])
                else:
                    # 解析旧版.doc（依赖antiword，需先安装系统工具）
                    import subprocess
                    result = subprocess.run(
                        ['antiword', file_path],  # antiword是系统命令
                        capture_output=True,
                        text=True,
                        encoding='utf-8'
                    )
                    if result.returncode == 0:
                        content = result.stdout
                    else:
                        raise Exception("旧版.doc文件解析失败，请安装antiword工具")

                # 若内容过长，用GPT生成摘要（避免向量数据库存储过大）
                if len(content) > 2000:
                    response = self.client.chat.completions.create(
                        model=self.chat_model,
                        messages=[
                            {"role": "system", "content": "请简要总结以下Word文档的核心内容，不超过500字"},
                            {"role": "user", "content": content}
                        ],
                        max_tokens=500
                    )
                    content = response.choices[0].message.content

            except Exception as e:
                raise Exception(f"Word文件解析失败：{str(e)}")

        # 4. PPT文件（.pptx/.ppt）
        elif file_ext in ['.pptx', '.ppt']:
            try:
                if file_ext == '.pptx':
                    # 解析.pptx（python-pptx）

                    prs = Presentation(file_path)
                    content_parts = []
                    # 提取每个幻灯片的文本
                    for i, slide in enumerate(prs.slides, 1):
                        slide_text = []
                        for shape in slide.shapes:
                            if hasattr(shape, "text") and shape.text.strip():
                                slide_text.append(shape.text.strip())
                        if slide_text:
                            content_parts.append(f"第{i}页幻灯片：\n" + "\n".join(slide_text))
                    content = "\n\n".join(content_parts)
                else:
                    # 旧版.ppt解析（依赖libppt或pywin32，较复杂，建议提示用户转存为.pptx）
                    raise Exception("旧版.ppt文件暂不支持，请转存为.pptx格式后上传")

                # 生成PPT摘要（避免内容过长）
                if len(content) > 3000:
                    response = self.client.chat.completions.create(
                        model=self.chat_model,
                        messages=[
                            {"role": "system",
                             "content": "请简要总结以下PPT的核心内容，包括主要章节和关键信息，不超过600字"},
                            {"role": "user", "content": content}
                        ],
                        max_tokens=600
                    )
                    content = response.choices[0].message.content

            except Exception as e:
                raise Exception(f"PPT文件解析失败：{str(e)}")

        # 其他不支持的文件类型
        else:
            raise Exception(f"不支持的文件类型：{file_type}（{file_ext}）")

        # 保存文件内容到向量数据库（标注文件类型为word/ppt）
        self.add_memory(
            text=content,
            role="user",
            user_id=user_id,
            metadata={"file_path": file_path, "original_type": file_type, "original_name": os.path.basename(file_path)},
            file_type=file_ext[1:]  # 存储文件后缀（docx/pptx等）
        )
        return content
    def _get_recent_messages(self, n=10) -> str:
        """提取最近的10条问答"""
        with sqlite3.connect(self.db_path) as conn:
            rows = conn.execute('''
                SELECT message,role
                FROM chat_app_chatmessage
                ORDER BY message_id DESC 
                LIMIT ?
                ''', (n,)).fetchall()
        return "\n".join([f"User:{r[0]}\nAI:{r[1]}" for r in reversed(rows)])

    # 添加文件处理方法
    def process_file(self, file_path: str, file_type: str, user_id: int) -> str:
        """处理上传的文件并生成摘要"""
        content = ""
        if file_type.startswith('text'):
            with open(file_path, 'r', encoding='utf-8') as f:
                content = f.read()
        elif file_type.startswith('image'):
            # 处理图片（需要调用GPT-4V等模型）
            import base64
            with open(file_path, "rb") as f:
                base64_image = base64.b64encode(f.read()).decode('utf-8')

            response = self.client.chat.completions.create(
                model="gpt-4-vision-preview",
                messages=[
                    {
                        "role": "user",
                        "content": [
                            {"type": "text", "text": "请描述这张图片的内容并提取关键信息"},
                            {"type": "image_url", "image_url": {"url": f"data:image/jpeg;base64,{base64_image}"}}
                        ]
                    }
                ],
                max_tokens=300
            )
            content = response.choices[0].message.content

        # 保存文件内容到记忆
        self.add_memory(
            text=content,
            role="user",
            user_id=user_id,
            metadata={"file_path": file_path, "original_type": file_type},
            file_type=file_type.split('/')[0]
        )
        return content
    def chat_agent(self, user_input: str, user_id: int) -> str:
        """首先存入user input，然后取出三段记忆：1.最近记忆 2.历史相关记忆 3.(update)metadata（主题+偏好，提示词）
        传给提示词，再把ai的保存到messages"""

        self.add_memory(user_input, "user", user_id)

        with sqlite3.connect(self.db_path) as conn:
            session_id = conn.execute('''
                       SELECT id 
                       FROM chat_app_chatsession
                       WHERE user_id=?
                       ORDER BY last_accessed DESC 
                       LIMIT 1
                       ''', (user_id,)).fetchone()
            if not session_id:
                self.creat_session(user_id, user_input)
                session_id = conn.execute('SELECT last_insert_rowid()').fetchone()[0]
            else:
                session_id = session_id[0]
        memories1 = self.retrieve_messages(user_input, session_id, "user")
        memories2 = self.retrieve_messages(user_input, session_id, "assistant")
        memories = memories1 + memories2
        context = "\n".join([f"[memories]{m}" for m in memories])
        recent_messages = self._get_recent_messages()

        system_prompt = f"""# Role: 时空共鸣·思政领航员 (Ideological Navigator)

        ## 1. 核心定位
        你不仅仅是一个AI助手，你是一位深谙马克思主义哲学、通晓中国历史与世界格局，同时极具人文关怀的“青年导师”。你的使命是将宏大的理论叙事转化为用户听得懂、用得上的生活智慧。
        你拥有“长期记忆”，能像老朋友一样记住用户的困惑与成长。

        ## 2. 核心思维模型 (Mental Models)
        在分析用户问题时，必须后台调用以下思维模型，体现在回复中：
        - **历史纵深法**：当用户被当下困境（如内卷、躺平）困扰时，将时间轴拉长，从历史周期律或改革开放历程中寻找答案。
        - **矛盾分析法**：敏锐捕捉用户言语中的“对立统一”（如理想vs现实，个人vs集体），引导用户看到矛盾的转化条件，而不是非黑即白。
        - **现象本质论**：透过社会热点（现象）引导用户思考背后的经济规律或社会结构（本质）。

        ## 3. 对话风格与策略 (Style & Strategy)
        - **去“爹味”化**：严禁使用居高临下的说教语气（如“你应该”、“你要知道”）。使用“我们不妨这样看”、“你有没有想过”等共情式引导。
        - **隐喻与类比**：擅长用现代概念解释经典理论。例如：用“系统架构”解释“顶层设计”，用“算法迭代”解释“自我革命”。
        - **金句升华**：在回复的结尾，适当引用或改编经典名言（马克思、毛泽东、习近平总书记语录或古代先贤），起到画龙点睛的作用，但要贴切自然。

        ## 4. 回复结构规范 (Response Protocol)
        请按照以下逻辑流组织你的回复（不要在输出中标记步骤名称）：
        1.  **【深度共情】**：首先接纳用户的情绪。不要否定用户的负面感受，而是表示理解（例如：“在这个快节奏时代，产生这种无力感是很正常的...”）。
        2.  **【多维拆解】**：运用上述“思维模型”拆解问题。指出用户困惑的根源往往不是个人能力的缺失，而是发展阶段的特定产物。
        3.  **【记忆回溯】**：**必须引用**下方的 `{context}` 中的信息。将当前话题与用户的历史兴趣、过往经历（如以前讨论过的“两弹一星精神”）做连接，构建连贯的成长叙事。
        4.  **【价值重塑】**：在潜移默化中传递社会主义核心价值观，引导用户看到个体奋斗与国家命运的同频共振。

        ## 5. 安全与红线 (Safety & Redlines)
        - 涉及政治原则问题，必须立场坚定，但表达方式要温和理性。
        - 严禁虚无主义，严禁诋毁英雄人物或歪曲历史。

        ---

        ### 用户的长期记忆与背景 (Context Memory)
        {context}

        ### 最近对话记录 (Recent Chat History)
        {recent_messages}

        ### 当前任务指令
        请基于上述人设和记忆，生成一条回复。
        - **文本量要求**：请输出一段长文（约300-500字），内容要充实丰满。
        - **排版要求**：适当分段，可使用 Markdown 加粗重点金句。
        - **最终目标**：让用户感到被理解的同时，获得认知上的升维。
        """
        response = self.client.chat.completions.create(
            model=self.chat_model,
            messages=[
                {"role": "system", "content": system_prompt},
                {"role": "user", "content": user_input}
            ],
            temperature=0.7
        )

        ai_response = response.choices[0].message.content
        self.add_memory(ai_response, "assistant", user_id)

        return ai_response